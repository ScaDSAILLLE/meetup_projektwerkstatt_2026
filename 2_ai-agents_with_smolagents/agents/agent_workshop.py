"""Workshop-Generator Agent using SmolAgents.

Dieser Agent generiert Workshop-Materialien: Markdown-Texte,
PowerPoint-Präsentationen und Python-Code-Beispiele.

Er nutzt einen ``ToolCallingAgent`` mit drei Tools:

1. **create_presentation** -- Erstellt eine PowerPoint-Präsentation aus
   einer Markdown-Datei. Slides werden anhand ``## Überschrift`` erstellt.

2. **write_file** -- Speichert Text als Datei ab (z.B. Markdown oder
   Python-Code).

3. **check_code** -- Überprüft Python-Code auf Lauffähigkeit.

Einstiegspunkt:

* ``create_agent()`` -- Factory-Funktion, die einen konfigurierten
  Agent zurückgibt::

      from agents.agent_workshop import create_agent

      agent = create_agent()
      answer = agent.run("Erstelle einen Workshop zum Thema KI")
"""

from __future__ import annotations

import subprocess as sp
import sys
from pathlib import Path

from pptx import Presentation
from smolagents import OpenAIServerModel, ToolCallingAgent, tool

try:
    import config
except ImportError:
    from pathlib import Path as _P

    sys.path.insert(0, str(_P(__file__).resolve().parents[1]))
    import config

OUTPUT_DIR: Path = config.OUTPUT_DIR


def _safe_path(name: str) -> Path:
    """Return a safe path inside ``OUTPUT_DIR``, creating the dir if needed."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR / Path(name).name


@tool
def create_presentation(quelldatei: str, name: str) -> str:
    """erstellt eine PowerPoint Präsentation für ein bestimmtes Thema, welches als Markdown vorliegt. Slides werden anhand ##Überschrift erstellt.

    Args:
        quelldatei: Dateiname der Markdown-Datei, z.B 'xx.md'
        name: Dateiname der Präsentation, z.B 'xx.pptx'
    """
    pfad = OUTPUT_DIR / Path(quelldatei).name
    quelle = pfad.read_text(encoding="utf-8")

    prs = Presentation()
    titel_inhalt_layout = prs.slide_layouts[1]

    welcome = prs.slides.add_slide(prs.slide_layouts[0])
    welcome.shapes.title.text = quelle.split("## ")[0].strip().lstrip("# ")
    welcome.placeholders[1].text = "Herzlich willkommen"

    for text in quelle.split("## ")[1:]:
        zeilen = text.strip().splitlines()
        slide = prs.slides.add_slide(titel_inhalt_layout)
        slide.shapes.title.text = zeilen[0]
        slide.placeholders[1].text = "\n".join(zeilen[1:6])

    ziel = _safe_path(name)
    prs.save(ziel)
    return f"Präsentation wurde gespeichert: {ziel}"


@tool
def write_file(name: str, text: str) -> str:
    """Speichert Text als Datei und mit einem Namen ab.

    Args:
        name: Bezeichnung der Datei mit Endung, z.B. 'vorlesung.md'
        text: Inhalt der Datei
    """
    ziel = _safe_path(name)
    ziel.write_text(text, encoding="utf-8")
    return f"Datei gespeichert: {ziel}"


@tool
def check_code(codename: str) -> str:
    """Überprüft Python Code eines generierten Beispiels auf Lauffähigkeit.

    Args:
        codename: Dateiname des Python-Code-Beispiels
    """
    pfad = OUTPUT_DIR / Path(codename).name
    if not pfad.is_file():
        return f"Datei nicht gefunden: {codename}"
    try:
        result = sp.run([sys.executable, str(pfad)], capture_output=True, text=True)
    except Exception as e:
        return f"Fehler bei der Ausführung: {e}"
    if result.returncode == 0:
        return "Code-Beispiel läuft erfolgreich."
    return f"Fehler:\n{result.stderr}"


def _require_api_key() -> str:
    key = config.SCADSAI_API_KEY
    if not key:
        raise RuntimeError(
            "SCADSAI_API_KEY is not set. "
            "Copy .env.example to .env and fill in SCADSAI_API_KEY=..."
        )
    return key


def create_agent() -> ToolCallingAgent:
    """Create and return a configured Workshop-Generator Agent.

    The agent uses the ScaDS.AI chat model for reasoning and tool selection.
    It has access to three tools:

    1. ``create_presentation`` -- creates a PPTX from a Markdown file
    2. ``write_file`` -- saves text as a file
    3. ``check_code`` -- checks if Python code runs without errors

    Returns:
        A configured ``ToolCallingAgent`` instance.

    Raises:
        RuntimeError: If ``SCADSAI_API_KEY`` is not set.
    """
    model = OpenAIServerModel(
        model_id=config.SCADSAI_CHAT_MODEL,
        api_base=config.SCADSAI_API_BASE,
        api_key=_require_api_key(),
    )
    # geändert, OW:
    # Kein globales tool_choice am Model setzen. SmolAgents reicht Tool-Optionen
    # nur bei echten Tool-Aufrufen weiter; das verhindert API-Fehler bei finalen
    # Textantworten ohne tools-Payload.

    return ToolCallingAgent(
        tools=[create_presentation, write_file, check_code],
        model=model,
        name="workshop_generator",
        description=(
            "Ein Agent der Workshop-Materialien generiert: Markdown-Texte, "
            "PowerPoint-Präsentationen und Python-Code-Beispiele."
        ),
        instructions=(
            "Du bist ein Workshop-Generator-Agent. Deine Aufgabe ist es, "
            "Workshop-Materialien zu generieren.\n\n"
            "Typischer Workflow:\n"
            "1. write_file: Schreibe einen Vorlesungstext in Markdown und "
            "speichere ihn.\n"
            "2. create_presentation: Erstelle aus dem Markdown eine "
            "PowerPoint-Präsentation.\n"
            "3. write_file: Schreibe ein Python-Code-Beispiel zum Thema.\n"
            "4. check_code: Überprüfe, ob der Code lauffähig ist.\n\n"
            "Gib am Ende eine Zusammenfassung der erstellten Dateien zurück."
        ),
        max_steps=10,
    )


__all__ = ["create_agent", "create_presentation", "write_file", "check_code"]


if __name__ == "__main__":
    agent = create_agent()
    response = agent.run(
        "1. Schreibe einen kurzen Vorlesungstext zu 'Agenten' in Markdown "
        "(eine '# ' Überschrift, drei '## ' Abschnitte) und speichere ihn mit "
        "write_file als 'vorlesung.md'. "
        "2. Rufe dann create_presentation(quelldatei='vorlesung.md', "
        "name='workshop.pptx') auf."
    )
    print(response)
